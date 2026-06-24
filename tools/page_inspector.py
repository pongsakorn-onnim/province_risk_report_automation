#!/usr/bin/env python3
"""
Universal Page Inspector
Reads ALL objects on any PowerPoint page including images, textboxes, and every formatting detail
Can inspect any page number, all pages, or slide masters in a presentation
"""

import os
import sys
# บังคับ encoding เป็น utf-8 เพื่อรองรับภาษาไทย
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR

# --- CLASS สำหรับ Redirect Output ลงไฟล์ ---
class Logger(object):
    def __init__(self, filename):
        self.terminal = sys.stdout
        self.log = open(filename, "w", encoding='utf-8')

    def write(self, message):
        # self.terminal.write(message) # ปิดไว้เพื่อให้หน้าจอไม่รก (ลงไฟล์อย่างเดียว)
        self.log.write(message)

    def flush(self):
        pass

# --- HELPER FUNCTIONS ---

def get_shape_type_name(shape_type):
    """Get human-readable shape type name"""
    shape_types = {}
    shape_type_mappings = [
        ('AUTO_SHAPE', "Auto Shape"), ('CALLOUT', "Callout"), ('CHART', "Chart"),
        ('CONNECTOR', "Connector"), ('FREEFORM', "Freeform"), ('GROUP', "Group"),
        ('LINE', "Line"), ('MEDIA', "Media"), ('OLE_CONTROL_OBJECT', "OLE Control Object"),
        ('OLE_OBJECT', "OLE Object"), ('PICTURE', "Picture"), ('PLACEHOLDER', "Placeholder"),
        ('TABLE', "Table"), ('TEXT_BOX', "Text Box"),
    ]
    for attr_name, display_name in shape_type_mappings:
        if hasattr(MSO_SHAPE_TYPE, attr_name):
            shape_types[getattr(MSO_SHAPE_TYPE, attr_name)] = display_name
    return shape_types.get(shape_type, f"Unknown ({shape_type})")

def get_alignment_name(alignment):
    """Get human-readable alignment name"""
    alignments = {}
    alignment_mappings = [
        ('LEFT', "Left"), ('CENTER', "Center"), ('RIGHT', "Right"),
        ('JUSTIFY', "Justify"), ('DISTRIBUTE', "Distribute"),
    ]
    for attr_name, display_name in alignment_mappings:
        if hasattr(PP_ALIGN, attr_name):
            alignments[getattr(PP_ALIGN, attr_name)] = display_name
    return alignments.get(alignment, f"Unknown ({alignment})")

def get_theme_color_name(theme_color):
    """Get human-readable theme color name"""
    theme_colors = {
        MSO_THEME_COLOR.ACCENT_1: "Accent 1", MSO_THEME_COLOR.ACCENT_2: "Accent 2",
        MSO_THEME_COLOR.ACCENT_3: "Accent 3", MSO_THEME_COLOR.ACCENT_4: "Accent 4",
        MSO_THEME_COLOR.ACCENT_5: "Accent 5", MSO_THEME_COLOR.ACCENT_6: "Accent 6",
        MSO_THEME_COLOR.BACKGROUND_1: "Background 1", MSO_THEME_COLOR.BACKGROUND_2: "Background 2",
        MSO_THEME_COLOR.DARK_1: "Dark 1", MSO_THEME_COLOR.DARK_2: "Dark 2",
        MSO_THEME_COLOR.LIGHT_1: "Light 1", MSO_THEME_COLOR.LIGHT_2: "Light 2",
        MSO_THEME_COLOR.HYPERLINK: "Hyperlink", MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "Followed Hyperlink"
    }
    return theme_colors.get(theme_color, f"Unknown ({theme_color})")

def extract_color_info(color):
    """Extract comprehensive color information"""
    color_info = {'type': type(color).__name__, 'rgb': None, 'theme_color': None, 'brightness': None, 'saturation': None, 'luminance': None}
    try:
        if hasattr(color, 'rgb') and color.rgb:
            color_info['rgb'] = color.rgb
            if hasattr(color.rgb, '_color_val'):
                r = (color.rgb._color_val >> 16) & 0xFF
                g = (color.rgb._color_val >> 8) & 0xFF
                b = color.rgb._color_val & 0xFF
                color_info['rgb_hex'] = f"#{r:02X}{g:02X}{b:02X}"
                color_info['rgb_values'] = f"RGB({r}, {g}, {b})"
        if hasattr(color, 'theme_color'):
            color_info['theme_color'] = get_theme_color_name(color.theme_color)
        if hasattr(color, 'brightness'):
            color_info['brightness'] = color.brightness
        if hasattr(color, 'saturation'):
            color_info['saturation'] = color.saturation
        if hasattr(color, 'luminance'):
            color_info['luminance'] = color.luminance
    except Exception as e:
        color_info['error'] = str(e)
    return color_info

def analyze_text_frame(text_frame, shape_index):
    """Analyze text frame in complete detail"""
    print(f"       TEXT FRAME ANALYSIS:")
    print(f"           Auto Size: {text_frame.auto_size}")
    print(f"           Margins: Left={text_frame.margin_left}, Right={text_frame.margin_right}, Top={text_frame.margin_top}, Bottom={text_frame.margin_bottom}")
    print(f"           Vertical Anchor: {text_frame.vertical_anchor}")
    print(f"           Word Wrap: {text_frame.word_wrap}")
    print(f"           Paragraphs: {len(text_frame.paragraphs)}")
    
    for para_idx, para in enumerate(text_frame.paragraphs):
        print(f"           PARAGRAPH {para_idx + 1}:")
        print(f"               Alignment: {get_alignment_name(para.alignment)}")
        print(f"               Space After: {para.space_after}")
        print(f"               Space Before: {para.space_before}")
        print(f"               Line Spacing: {para.line_spacing}")
        print(f"               Level: {para.level}")
        print(f"               Runs: {len(para.runs)}")
        
        if hasattr(para, 'bullet') and para.bullet:
            print(f"               Bullet: {para.bullet}")
        
        for run_idx, run in enumerate(para.runs):
            print(f'               RUN {run_idx + 1}: "{run.text}"' )
            
            font = run.font
            print(f"                   Font Name: {font.name}")
            print(f"                   Font Size: {font.size}")
            print(f"                   Bold: {font.bold}")
            print(f"                   Italic: {font.italic}")
            print(f"                   Underline: {font.underline}")
            
            font_attrs = ['strike', 'double_strike', 'subscript', 'superscript', 'small_caps', 'all_caps']
            for attr in font_attrs:
                if hasattr(font, attr):
                    value = getattr(font, attr)
                    print(f"                   {attr.replace('_', ' ').title()}: {value}")
                else:
                    print(f"                   {attr.replace('_', ' ').title()}: Not available")
            
            color_info = extract_color_info(font.color)
            print(f"                   Color Type: {color_info['type']}")
            if color_info.get('rgb_hex'):
                print(f"                   Color RGB: {color_info.get('rgb_values', 'N/A')} ({color_info['rgb_hex']})")
            if color_info.get('theme_color'):
                print(f"                   Theme Color: {color_info['theme_color']}")
            if color_info.get('brightness') is not None:
                print(f"                   Brightness: {color_info['brightness']}")
            if color_info.get('saturation') is not None:
                print(f"                   Saturation: {color_info['saturation']}")
            if color_info.get('luminance') is not None:
                print(f"                   Luminance: {color_info['luminance']}")
            if 'error' in color_info:
                print(f"                   Color Error: {color_info['error']}")
            
            if hasattr(font, 'language_id'):
                print(f"                   Language ID: {font.language_id}")
            if hasattr(font, 'rtl'):
                print(f"                   RTL: {font.rtl}")
            if hasattr(font, 'cs'):
                print(f"                   Complex Script: {font.cs}")

def analyze_picture(shape, shape_index):
    """Analyze picture in complete detail"""
    print(f"       PICTURE ANALYSIS:")
    print(f"           Image Type: {type(shape.image).__name__}")
    print(f"           Image Size: {shape.image.size}")
    print(f"           Image Format: {shape.image.content_type}")
    
    if hasattr(shape.image, 'filename'):
        print(f"           Filename: {shape.image.filename}")
    if hasattr(shape.image, 'path'):
        print(f"           Path: {shape.image.path}")
    
    if hasattr(shape, 'crop_left'):
        print(f"           Crop - Left: {shape.crop_left}, Right: {shape.crop_right}, Top: {shape.crop_top}, Bottom: {shape.crop_bottom}")
    
    if hasattr(shape, 'brightness'):
        print(f"           Brightness: {shape.brightness}")
    if hasattr(shape, 'contrast'):
        print(f"           Contrast: {shape.contrast}")

def analyze_group(shape, shape_index, depth=0):
    """Analyze group in complete detail with nested group support"""
    indent = "    " * (depth + 1)
    print(f"       GROUP ANALYSIS (Depth {depth}):")
    print(f"           Child Shapes: {len(shape.shapes)}")
    
    for child_idx, child_shape in enumerate(shape.shapes):
        print(f"           CHILD {child_idx + 1}: {get_shape_type_name(child_shape.shape_type)}")
        if hasattr(child_shape, 'left'):
            print(f'               Position: ({child_shape.left.inches:.3f}", {child_shape.top.inches:.3f}")')
        if hasattr(child_shape, 'width'):
            print(f'               Size: {child_shape.width.inches:.2f}" x {child_shape.height.inches:.2f}" ')
        
        if hasattr(child_shape, 'text_frame') and child_shape.text_frame and child_shape.text.strip():
            print(f'               Text Content: "{child_shape.text.strip()}"' )

        if child_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            analyze_group(child_shape, child_idx + 1, depth + 1)

def analyze_table(shape, shape_index):
    """Analyze table in complete detail"""
    print(f"       TABLE ANALYSIS:")
    print(f"           Rows: {len(shape.table.rows)}")
    print(f"           Columns: {len(shape.table.columns)}")
    
    for row_idx, row in enumerate(shape.table.rows):
        for col_idx, cell in enumerate(row.cells):
            if cell.text.strip():
                print(f'           Cell [{row_idx+1},{col_idx+1}]: "{cell.text.strip()}"' )

def analyze_chart(shape, shape_index):
    """Analyze chart in complete detail"""
    print(f"       CHART ANALYSIS:")
    print(f"           Chart Type: {type(shape.chart).__name__}")
    if hasattr(shape.chart, 'chart_type'):
        print(f"           Chart Type ID: {shape.chart.chart_type}")
    
    if hasattr(shape.chart, 'chart_data'):
        print(f"           Has Chart Data: Yes")
    if hasattr(shape.chart, 'has_title'):
        print(f"           Has Title: {shape.chart.has_title}")

def analyze_surface_shapes(surface):
    """Analyze all shapes on a given surface (slide or layout)."""
    print(f"   Total shapes: {len(surface.shapes)}")
    if hasattr(surface, 'slide_id'):
        print(f"   Slide ID: {surface.slide_id}")
    print()

    for i, shape in enumerate(surface.shapes):
        # [UPDATED] เพิ่มการ Print Shape Name เพื่อนำไปใช้ใน Config
        print(f"SHAPE {i+1}: '{shape.name}' ({get_shape_type_name(shape.shape_type)})")
        
        if hasattr(shape, 'left') and shape.left is not None:
            print(f'   Position: ({shape.left.inches:.3f}", {shape.top.inches:.3f}")')
        if hasattr(shape, 'width') and shape.width is not None:
            print(f'   Size: {shape.width.inches:.2f}" x {shape.height.inches:.2f}" ')
            if hasattr(shape, 'height') and shape.height is not None:
                 print(f"   Area: {shape.width.inches * shape.height.inches:.2f} sq inches")
        print(f"   Z-Order: {i} (0 = back, higher = front)")
        if hasattr(shape, 'rotation'):
            print(f"   Rotation: {shape.rotation}°")
        if hasattr(shape, 'fill'):
            fill = shape.fill
            print(f"   Fill Type: {type(fill).__name__}")
            try:
                if hasattr(fill, 'fore_color') and fill.fore_color is not None:
                    fill_color = extract_color_info(fill.fore_color)
                    print(f"   Fill Color: {fill_color['type']}")
                    # [UPDATED] ใช้ .get() เพื่อป้องกัน KeyError
                    if fill_color.get('rgb_hex'):
                        print(f"   Fill RGB: {fill_color['rgb_values']} ({fill_color['rgb_hex']})")
            except (TypeError, AttributeError) as e:
                print(f"   Fill Color: Not accessible ({type(e).__name__})")
        if hasattr(shape, 'line'):
            line = shape.line
            print(f"   Line Color: {extract_color_info(line.color)['type']}")
            if hasattr(line, 'width'):
                print(f"   Line Width: {line.width}")
        
        text_content = ""
        is_placeholder = False
        if hasattr(shape, 'text_frame') and shape.text_frame and shape.text.strip():
            text_content = shape.text.strip()
        elif hasattr(shape, 'text') and shape.text.strip():
             text_content = shape.text.strip()
             is_placeholder = True

        if text_content:
            if is_placeholder:
                print(f'   Text Content (Placeholder): "{text_content}"' )
            else:
                print(f'   Text Content: "{text_content}"' )
            if hasattr(shape, 'text_frame') and shape.text_frame:
                analyze_text_frame(shape.text_frame, i+1)

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            analyze_picture(shape, i+1)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            analyze_group(shape, i+1)
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            analyze_chart(shape, i+1)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            analyze_table(shape, i+1)
        print("   " + "=" * 80)
        print()

    print(f"SHAPE TYPE SUMMARY:")
    shape_counts = {}
    for shape in surface.shapes:
        shape_type = get_shape_type_name(shape.shape_type)
        shape_counts[shape_type] = shape_counts.get(shape_type, 0) + 1
    
    for shape_type, count in sorted(shape_counts.items()):
        print(f"   {shape_type}: {count}")
    
    text_shapes = [s for s in surface.shapes if (hasattr(s, 'text_frame') and s.text_frame and s.text.strip()) or (hasattr(s, 'text') and s.text.strip())]
    print(f"\nTEXT SUMMARY:")
    print(f"   Text shapes: {len(text_shapes)}")
    
    for i, shape in enumerate(text_shapes):
        text = shape.text.strip()
        preview = text[:50] + "..." if len(text) > 50 else text
        print(f'   {i+1}. "{preview}"' )
    
    image_shapes = [s for s in surface.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    print(f"\nIMAGE SUMMARY:")
    print(f"   Image shapes: {len(image_shapes)}")
    
    for i, shape in enumerate(image_shapes):
        print(f'   {i+1}. Size: {shape.width.inches:.2f}" x {shape.height.inches:.2f}" at ({shape.left.inches:.2f}", {shape.top.inches:.2f})')
    
    group_shapes = [s for s in surface.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]
    print(f"\nGROUP SUMMARY:")
    print(f"   Group shapes: {len(group_shapes)}")
    
    for i, shape in enumerate(group_shapes):
        print(f'   {i+1}. Size: {shape.width.inches:.2f}" x {shape.height.inches:.2f}" at ({shape.left.inches:.2f}", {shape.top.inches:.2f}") - {len(shape.shapes)} children')
    
    table_shapes = [s for s in surface.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
    print(f"\nTABLE SUMMARY:")
    print(f"   Table shapes: {len(table_shapes)}")
    
    for i, shape in enumerate(table_shapes):
        print(f'   {i+1}. Size: {shape.width.inches:.2f}" x {shape.height.inches:.2f}" at ({shape.left.inches:.2f}", {shape.top.inches:.2f}") - {len(shape.table.rows)}x{len(shape.table.columns)}')
    
    chart_shapes = [s for s in surface.shapes if s.shape_type == MSO_SHAPE_TYPE.CHART]
    print(f"\nCHART SUMMARY:")
    print(f"   Chart shapes: {len(chart_shapes)}")
    
    for i, shape in enumerate(chart_shapes):
        print(f'   {i+1}. Size: {shape.width.inches:.2f}" x {shape.height.inches:.2f}" at ({shape.left.inches:.2f}", {shape.top.inches:.2f}") - {type(shape.chart).__name__}')

def analyze_slide(slide, page_number):
    """Analyze a single slide comprehensively"""
    print(f"PAGE {page_number} COMPLETE ANALYSIS:")
    analyze_surface_shapes(slide)

def analyze_layout(layout, layout_name):
    """Analyze a single slide layout comprehensively"""
    print(f"LAYOUT '{layout_name}' COMPLETE ANALYSIS:")
    analyze_surface_shapes(layout)

def inspect_footers(prs):
    """Find and perform a detailed analysis of footer shapes in slide masters."""
    print("INSPECTING FOOTERS IN SLIDE MASTERS...")
    found_footers = 0
    for i, master in enumerate(prs.slide_masters):
        print(f"\n--- Master {i+1} ---")
        for j, layout in enumerate(master.slide_layouts):
            layout_name = layout.name if hasattr(layout, 'name') else f'Layout_{j}'
            for shape_idx, shape in enumerate(layout.shapes):
                if not shape.has_text_frame:
                    continue
                
                text = shape.text.strip()
                # Keywords to identify a footer shape
                if '‹#›' in text or '©' in text or 'สถานการณ์น้ำ' in text:
                    found_footers += 1
                    print(f"\n{'-'*40}")
                    print(f"  FOUND FOOTER in Layout: '{layout_name}' (Shape {shape_idx + 1})")
                    print(f"  Shape Type: {get_shape_type_name(shape.shape_type)}")
                    print(f'  Text: "{text}"')
                    analyze_text_frame(shape.text_frame, shape_idx + 1)
                    print(f"{'-'*40}\n")

    if found_footers == 0:
        print("  No potential footer shapes found in any slide masters.")
    else:
        print(f"SUCCESS: Found and analyzed {found_footers} potential footer(s).")

def inspect_masters(prs):
    """Perform a comprehensive analysis of all slide masters and layouts."""
    print(f"INSPECTING ALL SLIDE MASTERS AND LAYOUTS...")
    print(f"\nPresentation has {len(prs.slide_masters)} slide masters")
    for i, master in enumerate(prs.slide_masters):
        print(f"\n{'='*100}")
        print(f"MASTER {i+1} ANALYSIS")
        print(f"{ '='*100}")
        for j, layout in enumerate(master.slide_layouts):
            layout_name = layout.name if hasattr(layout, 'name') else f'Layout_{j}'
            print(f"\n{'-'*100}")
            analyze_layout(layout, layout_name)
    print(f"\nSUCCESS: Comprehensive master/layout inspection completed.")

def inspect_slides(prs, page_number=None):
    """Inspect all slides or a single slide."""
    total_slides = len(prs.slides)
    print(f"INSPECTING SLIDES...")
    print(f"\nPresentation has {total_slides} slides")
    if page_number:
        if page_number < 1 or page_number > total_slides:
            print(f"ERROR: Page {page_number} not found. Available pages: 1-{total_slides}")
            return
        slide = prs.slides[page_number - 1]
        analyze_slide(slide, page_number)
        print(f"\nSUCCESS: Inspection of page {page_number} completed.")
        print(f"Total objects analyzed: {len(slide.shapes)}")
    else:
        for page_num in range(1, total_slides + 1):
            slide = prs.slides[page_num - 1]
            print(f"\n{'='*100}")
            print(f"PAGE {page_num} ANALYSIS")
            print(f"{ '='*100}")
            analyze_slide(slide, page_num)
        print(f"\nSUCCESS: Comprehensive inspection of all {total_slides} slides completed.")

def get_latest_result_file():
    """Find the most recently modified file in the output directory."""
    import glob
    result_patterns = [
        "output/drought/*.pptx", "output/flood/*.pptx",
        "output/**/*.pptx"
    ]
    possible_paths = []
    for pattern in result_patterns:
        found_files = glob.glob(pattern, recursive=True)
        valid_files = [f for f in found_files if not os.path.basename(f).startswith('~$')]
        possible_paths.extend(valid_files)
    
    if not possible_paths:
        return None
        
    possible_paths.sort(key=os.path.getmtime, reverse=True)
    return possible_paths[0]

def print_help():
    """Prints the help message for the script."""
    print("Universal Page Inspector")
    print("Usage: python page_inspector.py [target_file] [what_to_inspect]")
    print("\n--- 1. Target File (Which presentation to open?)")
    print("   (no file path)      Inspects the latest result file in the output/ directory.")
    print("   --template [name]     Inspects a named template (e.g., drought_v2.pptx).")
    print("   <path/to/file.pptx>   Inspects any .pptx file by its path.")
    print("\n--- 2. Inspection Mode (What to look at inside the file?)")
    print("   (no mode)             Default. Analyzes every regular slide.")
    print("   <page_number>         Analyzes only the specified slide number.")
    print("   --masters             Dumps all details for every shape on all slide masters.")
    print("   --footers             Finds and performs a detailed analysis on only footer shapes.")
    print("   --help                Displays this help message.")

if __name__ == "__main__":
    args = sys.argv[1:]

    if "--help" in args or "-h" in args:
        print_help()
        sys.exit(0)

    target_file = None
    page_number = None
    inspection_mode = 'slides' 

    if "--masters" in args:
        inspection_mode = 'masters'
        args.remove("--masters")
    
    if "--footers" in args:
        inspection_mode = 'footers'
        args.remove("--footers")

    if "--template" in args:
        try:
            template_arg_index = args.index("--template")
            if template_arg_index + 1 < len(args) and not args[template_arg_index + 1].startswith('--'):
                template_name = args.pop(template_arg_index + 1)
                # [UPDATED] ใช้ folder 'templates' (เติม s) ตาม V2 Structure
                target_file = os.path.join('templates', template_name)
            else:
                target_file = os.path.join('templates', "drought_v2.pptx")
            args.pop(template_arg_index)
        except ValueError:
            pass

    if args:
        arg = args[0]
        if arg.endswith('.pptx'):
            target_file = arg
        elif arg.isdigit():
            page_number = int(arg)

    if target_file is None:
        target_file = get_latest_result_file()
    
    if not target_file or not os.path.exists(target_file):
        print(f"ERROR: Target file not found: {target_file}")
        sys.exit(1)

    # --- Output Setup (Save to File Logic) ---
    if not os.path.exists('inspection_pages'):
        os.makedirs('inspection_pages')

    filename = os.path.basename(target_file).replace('.pptx', '')
    if inspection_mode == 'slides' and page_number:
        out_name = f"{filename}_Page{page_number}.txt"
    else:
        out_name = f"{filename}_{inspection_mode.title()}.txt"
    
    output_path = os.path.join('inspection_pages', out_name)

    print("=" * 60)
    print(f"Running Inspection...")
    print(f"Target: {target_file}")
    print(f"Output will be saved to: {output_path}")
    print("=" * 60)

    # Redirect stdout ไปลงไฟล์ด้วย Logger class
    sys.stdout = Logger(output_path)

    try:
        prs = Presentation(target_file)
        if inspection_mode == 'masters':
            inspect_masters(prs)
        elif inspection_mode == 'footers':
            inspect_footers(prs)
        else: 
            inspect_slides(prs, page_number)
    except Exception as e:
        # ถ้า Error ก็จะลงไฟล์ด้วยเช่นกัน
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()