import logging
from copy import deepcopy
from lxml import etree
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

_BLACK = RGBColor(0x00, 0x00, 0x00)
_RED   = RGBColor(0xC0, 0x00, 0x00)
_NO_RISK_TEXT = "ไม่มีพื้นที่เสี่ยง"


def _get_template_rPr(cell):
    """Return a deep copy of the first run's rPr element found in the cell, or None."""
    txBody = cell._tc.find(qn('a:txBody'))
    if txBody is None:
        return None
    for p_elem in txBody.findall(qn('a:p')):
        for r_elem in p_elem.findall(qn('a:r')):
            rPr = r_elem.find(qn('a:rPr'))
            if rPr is not None:
                return deepcopy(rPr)
    return None


def _get_table_font_rPr(table):
    """Scan all rows (including header) for any rPr with an explicit <a:latin> font element."""
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            rPr = _get_template_rPr(table.cell(r, c))
            if rPr is not None and rPr.find(qn('a:latin')) is not None:
                return rPr
    return None


def _get_no_risk_rPr(table, col, font_rPr):
    """Return the rPr to use for a no-risk merged cell (always red + bold).

    Borrows font/size from the template cell's rPr if available so the DB Heavent
    font is preserved, but always forces explicit red RGB regardless of template color.
    """
    cell = table.cell(1, col)
    if cell.is_merge_origin:
        rPr = _get_template_rPr(cell)
        if rPr is not None:
            if rPr.find(qn('a:latin')) is None and font_rPr is not None:
                latin = font_rPr.find(qn('a:latin'))
                if latin is not None:
                    rPr.insert(0, deepcopy(latin))
            return _make_rPr(rPr, bold=True, color=_RED)
    return _make_rPr(font_rPr, bold=True, color=_RED)


def _make_rPr(template_rPr, bold, color=None):
    """Build an rPr element from a template, overriding bold and optionally color."""
    if template_rPr is not None:
        rPr = deepcopy(template_rPr)
    else:
        rPr = etree.Element(qn('a:rPr'))

    rPr.set('b', '1' if bold else '0')

    if color is not None:
        for sf in rPr.findall(qn('a:solidFill')):
            rPr.remove(sf)
        sf = etree.Element(qn('a:solidFill'))
        srgb = etree.SubElement(sf, qn('a:srgbClr'))
        srgb.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')
        # Insert after a:ln (if present) and before effectLst/latin/etc.
        # OOXML rPr schema: ln → fill → effect → highlight → uLn → uFill → latin …
        ln = rPr.find(qn('a:ln'))
        pos = list(rPr).index(ln) + 1 if ln is not None else 0
        rPr.insert(pos, sf)

    return rPr


def _attach_rPr(run, rPr):
    """Replace the run element's rPr with the given one (must be first child)."""
    r_elem = run._r
    old = r_elem.find(qn('a:rPr'))
    if old is not None:
        r_elem.remove(old)
    r_elem.insert(0, rPr)


def _split_if_merged(cell):
    if cell.is_merge_origin:
        cell.split()


def _normalize_col(table, col):
    """Split any merged cells in data rows 1-3 of the given column."""
    for row in range(1, 4):
        _split_if_merged(table.cell(row, col))


def _write_no_risk_merged(table, col, no_risk_rPr):
    """Merge rows 1-3 and write the no-risk label centered, preserving template color."""
    from pptx.enum.text import MSO_ANCHOR
    table.cell(1, col).merge(table.cell(3, col))
    merged = table.cell(1, col)
    merged.text = _NO_RISK_TEXT
    merged.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = merged.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if p.runs:
        _attach_rPr(p.runs[0], deepcopy(no_risk_rPr))


def _write_plain_cell(cell, value, fallback_rPr=None):
    """Overwrite first run's text in place, preserving template font formatting."""
    _cell_rPr = _get_template_rPr(cell)
    template_rPr = _cell_rPr if _cell_rPr is not None else fallback_rPr
    tf = cell.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        r = p.runs[0]
        r.text = str(value)
        _attach_rPr(r, _make_rPr(template_rPr, bold=False, color=_BLACK))
        for extra in p.runs[1:]:
            extra.text = ""
        for extra_p in tf.paragraphs[1:]:
            for extra_r in extra_p.runs:
                extra_r.text = ""
    else:
        cell.text = str(value)  # resets paragraph — set alignment after
        if tf.paragraphs[0].runs:
            _attach_rPr(tf.paragraphs[0].runs[0],
                        _make_rPr(template_rPr, bold=False, color=_BLACK))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def _write_detail_cell(cell, detail_text, fallback_rPr=None):
    """Write amphoe:tambon lines — bold 'amphoe:' run + normal tambon run per line."""
    lines = [ln.strip() for ln in detail_text.replace('\r', '').split('\n') if ln.strip()]

    # Capture template rPr BEFORE clearing; fall back to table-level font rPr when
    # this cell is an empty span cell (from a split merge) with no runs of its own.
    _cell_rPr = _get_template_rPr(cell)
    template_rPr = _cell_rPr if _cell_rPr is not None else fallback_rPr

    cell.text = ""
    tf = cell.text_frame

    # Remove empty run left by cell.text = ""
    p0_elem = tf.paragraphs[0]._p
    for r_elem in p0_elem.findall(qn('a:r')):
        p0_elem.remove(r_elem)

    if not lines:
        return

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if ':' in line:
            amphoe_part, tambon_part = line.split(':', 1)
            r1 = p.add_run()
            r1.text = f"{amphoe_part}:"
            _attach_rPr(r1, _make_rPr(template_rPr, bold=True, color=_BLACK))

            r2 = p.add_run()
            r2.text = f" {tambon_part.strip()}"
            _attach_rPr(r2, _make_rPr(template_rPr, bold=False, color=_BLACK))
        else:
            r = p.add_run()
            r.text = line
            _attach_rPr(r, _make_rPr(template_rPr, bold=False, color=_BLACK))


def get_table(slide, shape_name: str):
    from src.ppt_tools.image_handler import _find_shape_recursive
    shape = _find_shape_recursive(slide.shapes, shape_name)
    if shape is None:
        logger.warning(f"Table shape '{shape_name}' not found")
        return None
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        logger.warning(f"Shape '{shape_name}' is not a table")
        return None
    return shape.table


def fill_province_6m_table(table, province_6m: dict):
    """Fill 4-row × 4-col province 6m summary table.

    Cols 1-3: เสี่ยงท่วม, เสี่ยงแล้ง, เสี่ยงทั้งท่วมทั้งแล้ง
    Rows 1-3: อำเภอ count, ตำบล count, detail text
    """
    flood  = province_6m.get("เสี่ยงท่วม", {})
    drought = province_6m.get("เสี่ยงแล้ง", {})
    both   = province_6m.get("เสี่ยงทั้งท่วมทั้งแล้ง", {})

    font_rPr = _get_table_font_rPr(table)
    for col_idx, data in enumerate([flood, drought, both], start=1):
        no_risk_rPr = _get_no_risk_rPr(table, col_idx, font_rPr)
        _normalize_col(table, col_idx)
        if not data.get("amphoe"):
            _write_no_risk_merged(table, col_idx, no_risk_rPr)
        else:
            _write_plain_cell(table.cell(1, col_idx), data["amphoe"], fallback_rPr=font_rPr)
            _write_plain_cell(table.cell(2, col_idx), data["tambon"], fallback_rPr=font_rPr)
            _write_detail_cell(table.cell(3, col_idx), data.get("detail", ""), fallback_rPr=font_rPr)


def fill_monthly_table(table, monthly_data: dict, month_year_key: str):
    """Fill 4-row × 3-col monthly table.

    Cols 1-2: เสี่ยงท่วม, เสี่ยงแล้ง
    Rows 1-3: อำเภอ count, ตำบล count, detail text
    month_year_key: 'MM-YYYY' e.g. '06-2026'
    """
    month = monthly_data.get(month_year_key, {})
    flood  = month.get("เสี่ยงท่วม", {})
    drought = month.get("เสี่ยงแล้ง", {})

    font_rPr = _get_table_font_rPr(table)
    for col_idx, data in enumerate([flood, drought], start=1):
        no_risk_rPr = _get_no_risk_rPr(table, col_idx, font_rPr)
        _normalize_col(table, col_idx)
        if not data.get("amphoe"):
            _write_no_risk_merged(table, col_idx, no_risk_rPr)
        else:
            _write_plain_cell(table.cell(1, col_idx), data["amphoe"], fallback_rPr=font_rPr)
            _write_plain_cell(table.cell(2, col_idx), data["tambon"], fallback_rPr=font_rPr)
            _write_detail_cell(table.cell(3, col_idx), data.get("detail", ""), fallback_rPr=font_rPr)


def _write_region_detail_cell(cell, region_details: list, region_display: dict = None,
                              fallback_rPr=None):
    """Write province list grouped by region: bold 'ภาคX: ' header then normal province text."""
    if region_display is None:
        region_display = {}

    _cell_rPr = _get_template_rPr(cell)
    template_rPr = _cell_rPr if _cell_rPr is not None else fallback_rPr

    cell.text = ""
    tf = cell.text_frame
    p0_elem = tf.paragraphs[0]._p
    for r_elem in p0_elem.findall(qn('a:r')):
        p0_elem.remove(r_elem)

    if not region_details:
        return

    for i, entry in enumerate(region_details):
        region_key = entry.get("region", "")
        provinces  = entry.get("provinces", "")
        label = region_display.get(region_key, region_key)

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if label:
            r1 = p.add_run()
            r1.text = f"{label}: "
            _attach_rPr(r1, _make_rPr(template_rPr, bold=True, color=_BLACK))

        if provinces:
            r2 = p.add_run()
            r2.text = provinces
            _attach_rPr(r2, _make_rPr(template_rPr, bold=False, color=_BLACK))


def fill_national_table(table, national_data: dict, region_display: dict = None):
    """Fill 5-row × 4-col national 6m summary table.

    Cols 1-3: flood, drought, both
    Rows 1-4: province count, amphoe count, tambon count, detail text
    """
    flood  = national_data.get("flood", {})
    drought = national_data.get("drought", {})
    both   = national_data.get("both", {})

    font_rPr = _get_table_font_rPr(table)
    count_keys = ["province", "amphoe", "tambon"]
    for row_idx, key in enumerate(count_keys, start=1):
        for col_idx, data in enumerate([flood, drought, both], start=1):
            _write_plain_cell(table.cell(row_idx, col_idx), data.get(key, ""), fallback_rPr=font_rPr)

    for col_idx, data in enumerate([flood, drought, both], start=1):
        _write_region_detail_cell(table.cell(4, col_idx), data.get("detail", []),
                                  region_display=region_display, fallback_rPr=font_rPr)
