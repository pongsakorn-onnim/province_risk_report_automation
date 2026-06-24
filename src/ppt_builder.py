import logging
from pathlib import Path
from datetime import datetime
from pptx import Presentation

from src.ppt_tools.image_handler import replace_image_by_name, fetch_url, _find_shape_recursive, _make_placeholder
from src.ppt_tools.text_handler import (
    replace_text_by_name, replace_text_containing, replace_in_runs,
    replace_in_shape_runs, THAI_MONTHS,
    get_thai_month, get_buddhist_year, get_next_months, format_month_range,
)
from src.path_resolver import get_image_paths, get_excel_paths
from src.data_reader import read_province_6m, read_province_monthly, read_national_6m
from src.table_filler import (
    get_table, fill_province_6m_table, fill_monthly_table, fill_national_table,
)

logger = logging.getLogger(__name__)

TITLE_PREFIX = "สรุปคาดการณ์พื้นที่เสี่ยงใน 6 เดือนข้างหน้า"


def _month_code_key(yyyymm_code: str) -> str:
    """Convert "202606" → "06-2026" (summary_1m format)."""
    return f"{yyyymm_code[4:6]}-{yyyymm_code[:4]}"


def _update_cover(slide, year: int, start_month: int, province: str, cfg: dict):
    months = get_next_months(year, start_month, 6)
    be_year = get_buddhist_year(year)
    day = datetime.today().day
    thai_start = get_thai_month(start_month)
    thai_end = months[-1]["thai_name"]

    replace_text_by_name(slide, cfg["shapes"]["cover"]["date_text"],
                         f"วันที่ {day} {thai_start} {be_year}")
    replace_text_by_name(slide, cfg["shapes"]["cover"]["period_text"],
                         f"เดือน{thai_start}-{thai_end} {be_year}")
    # Surgically replace province name in every run (preserves <a:br/> line break)
    replace_in_runs(slide, "น่าน", province)


def _update_national_summary(slide, paths: dict, cfg: dict):
    s = cfg["shapes"]["national_summary"]
    replace_image_by_name(slide, s["all_map"],      paths["nat_all_map"])
    replace_image_by_name(slide, s["flood_graph"],  paths["nat_flood_graph"])
    replace_image_by_name(slide, s["drought_graph"],paths["nat_drought_graph"])
    replace_image_by_name(slide, s["both_graph"],   paths["nat_both_graph"])


def _update_region_summary(slide, paths: dict, region: str, cfg: dict,
                           slide_height_emu: int = None):
    s = cfg["shapes"]["region_summary"]

    bar_shape_names = {s["flood_graph"], s["drought_graph"], s["both_graph"]}
    bar_lefts = [sh.left for sh in slide.shapes if sh.name in bar_shape_names]
    right_limit = (min(bar_lefts) - _MAP_MARGIN_EMU) if bar_lefts else None

    replace_image_by_name(slide, s["all_map"],      paths["reg_all_map"],
                          preserve_ratio=True, right_limit_emu=right_limit,
                          slide_height_emu=slide_height_emu)
    replace_image_by_name(slide, s["flood_graph"],  paths["reg_flood_graph"])
    replace_image_by_name(slide, s["drought_graph"],paths["reg_drought_graph"])
    replace_image_by_name(slide, s["both_graph"],   paths["reg_both_graph"])

    region_display = cfg.get("region_display", {}).get(region, f"ภาค{region}")
    replace_in_runs(slide, "ภาคเหนือ", region_display)


_MAP_MARGIN_EMU = 45_720  # 0.05" gap between province map and adjacent bar chart


def _update_province_6m_maps(slide, paths: dict, province: str, cfg: dict,
                             slide_height_emu: int = None):
    s = cfg["shapes"]["province_6m_maps"]

    # Dynamically expand the map box rightward up to just before the adjacent bar chart
    bar_shape_names = {s["flood_graph"], s["drought_graph"], s["both_graph"]}
    bar_lefts = [sh.left for sh in slide.shapes if sh.name in bar_shape_names]
    right_limit = (min(bar_lefts) - _MAP_MARGIN_EMU) if bar_lefts else None

    replace_image_by_name(slide, s["all_map"],      paths["prov_all_map"],
                          preserve_ratio=True, right_limit_emu=right_limit,
                          slide_height_emu=slide_height_emu)
    replace_image_by_name(slide, s["flood_graph"],  paths["prov_flood_graph"])
    replace_image_by_name(slide, s["drought_graph"],paths["prov_drought_graph"])
    replace_image_by_name(slide, s["both_graph"],   paths["prov_both_graph"])
    replace_in_runs(slide, "น่าน", province)


def _update_province_6m_table(slide, province_6m: dict, province: str, cfg: dict):
    replace_in_runs(slide, "น่าน", province)
    shape_name = cfg["shapes"]["province_6m_table"]
    table = get_table(slide, shape_name)
    if table:
        fill_province_6m_table(table, province_6m)


def _update_monthly_slide(slide, month_cfg: dict, map_path: Path,
                          province_monthly: dict, month_code: str, province: str,
                          slide_height_emu: int = None):
    m = int(month_code[4:6])
    y = int(month_code[:4])
    new_month = get_thai_month(m)
    be_year   = str(get_buddhist_year(y))

    # Replace each dynamic part run-by-run so per-run color/font is preserved.
    # Replacing all 12 month names covers any yyyymm period, not just 202606.
    title_replacements = {old_m: new_month for old_m in THAI_MONTHS}
    title_replacements["น่าน"] = province
    # Cover Buddhist year drift (±2) in case template and run period differ
    for delta in range(-2, 3):
        candidate = str(int(be_year) + delta)
        if candidate != be_year:
            title_replacements[candidate] = be_year
    replace_in_shape_runs(slide, month_cfg["title"], title_replacements)
    replace_image_by_name(slide, month_cfg["map"], map_path,
                          preserve_ratio=True, slide_height_emu=slide_height_emu)
    table = get_table(slide, month_cfg["table"])
    if table:
        fill_monthly_table(table, province_monthly, _month_code_key(month_code))


def _update_risk_forecast(slide, risk_cfg: dict, year: int, start_month: int):
    months = get_next_months(year, start_month, 6)
    title = f"{risk_cfg['title_prefix']}เดือน{format_month_range(months)}"
    replace_text_by_name(slide, risk_cfg["title"], title)
    yyyymm = f"{year}{start_month:02d}"
    for i, shape_name in enumerate(risk_cfg["maps"], start=1):
        url = risk_cfg["url_pattern"].format(yyyymm=yyyymm, m=i)
        stream = fetch_url(url)
        if stream is not None:
            replace_image_by_name(slide, shape_name, stream)
        else:
            logger.warning(f"Risk image unavailable: {url}")
            shape = _find_shape_recursive(slide.shapes, shape_name)
            if shape is not None:
                ph = _make_placeholder(shape.width, shape.height, url)
                replace_image_by_name(slide, shape_name, ph)


def _update_combined_map(slide, paths: dict, cfg: dict):
    c = cfg["shapes"]["combined_map"]
    replace_image_by_name(slide, c["flood_pic"],   paths["comb_flood_map"])
    replace_image_by_name(slide, c["drought_pic"], paths["comb_drought_map"])
    replace_image_by_name(slide, c["both_pic"],    paths["comb_both_map"])


def _update_national_table(slide, national_data: dict, cfg: dict):
    table = get_table(slide, cfg["shapes"]["national_table"])
    if table:
        fill_national_table(table, national_data,
                            region_display=cfg.get("region_display", {}))


def build_report(config: dict, province: str, region: str, data_dir: Path,
                 yyyymm: str, year: int, start_month: int, output_dir: Path) -> Path:
    """Fill the template with province data and save the output .pptx."""
    template_path = Path(config["template_path"])
    output_dir.mkdir(parents=True, exist_ok=True)

    logger.info(f"Building report: province={province}, region={region}, period={yyyymm}")

    paths = get_image_paths(data_dir, region, province, yyyymm)
    excel = get_excel_paths(data_dir, region, province, yyyymm)

    logger.info("Reading Excel data...")
    province_6m      = read_province_6m(excel["province"])
    province_monthly = read_province_monthly(excel["province"])
    national_data    = read_national_6m(excel["national"])

    logger.info("Loading template...")
    prs = Presentation(str(template_path))
    slides = list(prs.slides)

    # Generate the 6 month codes for this period
    month_codes = []
    for i in range(6):
        m = (start_month - 1 + i) % 12 + 1
        y = year + (start_month - 1 + i) // 12
        month_codes.append(f"{y}{m:02d}")

    cfg = config

    logger.info("Slide 1: cover")
    _update_cover(slides[0], year, start_month, province, cfg)

    logger.info("Slide 4: national summary maps")
    _update_national_summary(slides[3], paths, cfg)

    logger.info("Slide 5: region summary")
    _update_region_summary(slides[4], paths, region, cfg,
                           slide_height_emu=prs.slide_height)

    logger.info("Slide 6: province 6-month maps")
    _update_province_6m_maps(slides[5], paths, province, cfg,
                             slide_height_emu=prs.slide_height)

    logger.info("Slide 7: province 6-month table")
    _update_province_6m_table(slides[6], province_6m, province, cfg)

    monthly_cfgs = cfg["shapes"]["monthly"]
    for i, (mc_cfg, mc) in enumerate(zip(monthly_cfgs, month_codes)):
        logger.info(f"Slide {8+i}: monthly {mc}")
        _update_monthly_slide(
            slides[7 + i], mc_cfg,
            paths["prov_monthly_maps"][i],
            province_monthly, mc, province,
            slide_height_emu=prs.slide_height,
        )

    logger.info("Slide 15: flood risk forecast")
    _update_risk_forecast(slides[14], cfg["shapes"]["flood_risk"], year, start_month)

    logger.info("Slide 16: drought risk forecast")
    _update_risk_forecast(slides[15], cfg["shapes"]["drought_risk"], year, start_month)

    logger.info("Slide 17: combined national map")
    _update_combined_map(slides[16], paths, cfg)

    logger.info("Slide 18: national summary table")
    _update_national_table(slides[17], national_data, cfg)

    out_file = output_dir / f"{yyyymm}_{province}_FloodDrought_report.pptx"
    prs.save(str(out_file))
    logger.info(f"Saved: {out_file}")
    return out_file
