import io
import logging
import os
import tempfile
from pathlib import Path
from pptx.slide import Slide
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw
import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry

logger = logging.getLogger(__name__)

_EMU_PER_PX = 9525  # 96 DPI

_session = None

def _get_session() -> requests.Session:
    global _session
    if _session is None:
        _session = requests.Session()
        retry = Retry(total=3, backoff_factor=0.3, status_forcelist=[500, 502, 503, 504])
        adapter = HTTPAdapter(max_retries=retry)
        _session.mount("https://", adapter)
        _session.mount("http://", adapter)
    return _session


def fetch_url(url: str) -> io.BytesIO | None:
    """Download an image from a URL, return BytesIO or None on failure."""
    try:
        resp = _get_session().get(url, timeout=15)
        resp.raise_for_status()
        return io.BytesIO(resp.content)
    except requests.exceptions.HTTPError as e:
        logger.warning(f"URL not found: {url} (HTTP {e.response.status_code})")
    except Exception as e:
        logger.warning(f"URL fetch failed: {url} ({e})")
    return None


def _wrap_text(draw, text: str, font, max_px: int) -> list[str]:
    lines = []
    while text:
        lo, hi = 1, len(text)
        while lo < hi:
            mid = (lo + hi + 1) // 2
            if draw.textlength(text[:mid], font=font) <= max_px:
                lo = mid
            else:
                hi = mid - 1
        lines.append(text[:lo])
        text = text[lo:]
    return lines


def _make_placeholder(width_emu: int, height_emu: int, missing_path) -> io.BytesIO:
    px_w = max(60, round(width_emu / _EMU_PER_PX))
    px_h = max(40, round(height_emu / _EMU_PER_PX))
    img = Image.new("RGB", (px_w, px_h), color=(255, 220, 220))
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, px_w - 1, px_h - 1], outline=(180, 0, 0), width=3)
    try:
        from PIL import ImageFont
        font = ImageFont.load_default(size=10)
    except TypeError:
        from PIL import ImageFont
        font = ImageFont.load_default()
    max_text_w = px_w - 12
    lines = ["Image not found:"] + _wrap_text(draw, str(missing_path), font, max_text_w)
    line_h = 13
    y = 5
    for line in lines:
        if y + line_h > px_h - 4:
            break
        draw.text((6, y), line, fill=(140, 0, 0), font=font)
        y += line_h
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def _find_shape_recursive(shapes, shape_name: str):
    for shape in shapes:
        if shape.name == shape_name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = _find_shape_recursive(shape.shapes, shape_name)
            if found:
                return found
    return None


def _fit_within_box(img_path: Path, box_left: int, box_top: int,
                    box_w: int, box_h: int,
                    right_limit_emu: int = None,
                    slide_height_emu: int = None) -> tuple[int, int, int, int]:
    """Return (left, top, width, height) in EMU that fits img inside the box, preserving ratio.

    right_limit_emu: hard right boundary — box_w is set to right_limit - box_left,
                     expanding or clamping as needed to avoid overlap.
    slide_height_emu: center the result vertically on the full slide instead of within the box.
    """
    with Image.open(str(img_path)) as im:
        img_w, img_h = im.size  # pixels
    if right_limit_emu is not None:
        box_w = right_limit_emu - box_left
    scale = min(box_w / img_w, box_h / img_h)  # EMU per pixel
    new_w = round(img_w * scale)
    new_h = round(img_h * scale)
    new_left = box_left + (box_w - new_w) // 2  # center within available space
    if slide_height_emu is not None:
        new_top = (slide_height_emu - new_h) // 2
    else:
        new_top = box_top + (box_h - new_h) // 2
    return new_left, new_top, new_w, new_h


def replace_image_by_name(slide: Slide, shape_name: str, new_image,
                          preserve_ratio: bool = False,
                          right_limit_emu: int = None,
                          slide_height_emu: int = None) -> bool:
    """Replace a named picture shape.

    new_image: file path (str/Path) or a BytesIO stream (e.g. from fetch_url).
    """
    if not new_image:
        logger.warning(f"No image given for shape '{shape_name}'")
        return False

    is_stream = isinstance(new_image, io.BytesIO)
    path_obj = None if is_stream else Path(new_image)
    use_placeholder = not is_stream and not path_obj.exists()
    if use_placeholder:
        logger.error(f"Image not found: {new_image} (shape '{shape_name}') — inserting placeholder")

    target_shape = _find_shape_recursive(slide.shapes, shape_name)
    if not target_shape:
        logger.warning(f"Shape '{shape_name}' not found on slide {slide.slide_id}")
        return False
    if target_shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        logger.warning(f"Shape '{shape_name}' is not a picture (type {target_shape.shape_type})")
        return False

    left = target_shape.left
    top = target_shape.top
    width = target_shape.width
    height = target_shape.height

    if preserve_ratio and not use_placeholder and not is_stream:
        left, top, width, height = _fit_within_box(
            path_obj, left, top, width, height, right_limit_emu, slide_height_emu
        )

    target_shape.element.getparent().remove(target_shape.element)

    try:
        if use_placeholder:
            image_source = _make_placeholder(width, height, path_obj)
        elif is_stream:
            new_image.seek(0)
            image_source = new_image
        else:
            image_source = str(path_obj)
        pic = slide.shapes.add_picture(image_source, left, top, width, height)
        success = not use_placeholder
    except Exception as e:
        logger.error(f"Failed to insert image for '{shape_name}': {e}")
        try:
            pic = slide.shapes.add_picture(
                _make_placeholder(width, height, path_obj or shape_name),
                left, top, width, height,
            )
            success = False
        except Exception as e2:
            logger.error(f"Placeholder also failed for '{shape_name}': {e2}")
            return False

    sp = pic._element
    spTree = sp.getparent()
    spTree.remove(sp)
    spTree.insert(2, sp)
    return success


def replace_image_in_group(slide: Slide, group_name: str, pic_name: str, new_image_path) -> bool:
    """Replace an image inside a named group by swapping the image blob in-place.

    Does NOT move the picture element — it stays inside the group with its
    original position and size intact. Only the r:embed reference is updated.
    """
    from pptx.oxml.ns import qn as _qn

    group = None
    for shape in slide.shapes:
        if shape.name == group_name and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group = shape
            break
    if group is None:
        logger.warning(f"Group '{group_name}' not found on slide {slide.slide_id}")
        return False

    target = None
    for s in group.shapes:
        if s.name == pic_name and s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            target = s
            break
    if target is None:
        logger.warning(f"Picture '{pic_name}' not found in group '{group_name}'")
        return False

    path_obj = Path(new_image_path)
    use_placeholder = not path_obj.exists()
    if use_placeholder:
        logger.error(f"Image not found: {new_image_path} (group '{group_name}' / '{pic_name}') — inserting placeholder")

    blip = target._element.find('.//' + _qn('a:blip'))
    if blip is None:
        logger.error(f"No <a:blip> in '{pic_name}' inside group '{group_name}'")
        return False

    _EMBED = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
    if use_placeholder:
        ph_bytes = _make_placeholder(target.width, target.height, path_obj)
        tmp_fd, tmp_path = tempfile.mkstemp(suffix='.png')
        try:
            os.write(tmp_fd, ph_bytes.read())
            os.close(tmp_fd)
            _, new_rId = slide.part.get_or_add_image_part(tmp_path)
            blip.set(_EMBED, new_rId)
        finally:
            os.unlink(tmp_path)
        return False
    else:
        _, new_rId = slide.part.get_or_add_image_part(str(path_obj))
        blip.set(_EMBED, new_rId)
        return True
