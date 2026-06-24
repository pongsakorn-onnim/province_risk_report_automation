import logging
from pptx.slide import Slide
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

THAI_MONTHS = [
    "มกราคม", "กุมภาพันธ์", "มีนาคม", "เมษายน", "พฤษภาคม", "มิถุนายน",
    "กรกฎาคม", "สิงหาคม", "กันยายน", "ตุลาคม", "พฤศจิกายน", "ธันวาคม",
]


def get_thai_month(month_idx: int) -> str:
    if 1 <= month_idx <= 12:
        return THAI_MONTHS[month_idx - 1]
    return ""


def get_buddhist_year(year: int) -> int:
    return year + 543


def format_month_range(months: list[dict]) -> str:
    """Format a list of month dicts into 'มิถุนายน – พฤศจิกายน 2569' style.

    Cross-year: 'ธันวาคม 2568 – มกราคม 2569'
    """
    if not months:
        return ""
    start, end = months[0], months[-1]
    sep = " – "  # en dash
    if start["buddhist_year"] == end["buddhist_year"]:
        return f"{start['thai_name']}{sep}{end['thai_name']} {start['buddhist_year']}"
    return (f"{start['thai_name']} {start['buddhist_year']}"
            f"{sep}{end['thai_name']} {end['buddhist_year']}")


def get_next_months(start_year: int, start_month: int, n: int) -> list[dict]:
    results = []
    for i in range(n):
        m = (start_month + i - 1) % 12 + 1
        y = start_year + (start_month + i - 1) // 12
        results.append({
            "year": y,
            "month": m,
            "thai_name": get_thai_month(m),
            "buddhist_year": get_buddhist_year(y),
        })
    return results


def _find_shape_recursive(shapes, shape_name: str):
    for shape in shapes:
        if shape.name == shape_name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = _find_shape_recursive(shape.shapes, shape_name)
            if found:
                return found
    return None


def replace_text_by_name(slide: Slide, shape_name: str, new_text: str) -> bool:
    target = _find_shape_recursive(slide.shapes, shape_name)
    if not target:
        logger.warning(f"Text shape '{shape_name}' not found")
        return False
    if not target.has_text_frame:
        logger.warning(f"Shape '{shape_name}' has no text frame")
        return False

    tf = target.text_frame
    if tf.paragraphs:
        p0 = tf.paragraphs[0]
        if p0.runs:
            p0.runs[0].text = str(new_text)
            for run in p0.runs[1:]:
                run.text = ""
            for para in tf.paragraphs[1:]:
                for run in para.runs:
                    run.text = ""
        else:
            target.text = str(new_text)
    else:
        target.text = str(new_text)
    return True


def replace_text_containing(slide: Slide, search_text: str, new_text: str) -> bool:
    """Find the first shape whose text contains search_text and replace its text.

    Works at the paragraph level: replaces run[0] of the first paragraph that
    contains search_text, clearing subsequent runs in that paragraph.
    Other paragraphs are left intact (preserves <a:br/> structures).
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full_text = tf.text
        if search_text not in full_text:
            continue
        # Find the paragraph that contains search_text
        for para in tf.paragraphs:
            if search_text in para.text:
                if para.runs:
                    para.runs[0].text = str(new_text)
                    for run in para.runs[1:]:
                        run.text = ""
                return True
    logger.warning(f"No shape found containing text: {search_text!r}")
    return False


def replace_in_shape_runs(slide: Slide, shape_name: str, replacements: dict) -> bool:
    """Apply {old: new} replacements in every run of a named shape, preserving per-run formatting."""
    target = _find_shape_recursive(slide.shapes, shape_name)
    if not target or not target.has_text_frame:
        logger.warning(f"Shape '{shape_name}' not found or has no text frame")
        return False
    for para in target.text_frame.paragraphs:
        for run in para.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)
    return True


def replace_in_runs(slide: Slide, old_substring: str, new_substring: str) -> int:
    """Replace old_substring with new_substring in every run across the slide.

    Surgical replacement — preserves all other text and line-break elements.
    Returns count of substitutions made.
    """
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old_substring in run.text:
                    run.text = run.text.replace(old_substring, new_substring)
                    count += 1
    return count
